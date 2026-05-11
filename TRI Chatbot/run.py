import base64
import hashlib
import json
import os
import tempfile
import time
import uuid
from pathlib import Path

from flask import Flask, Response, jsonify, request, send_from_directory, session
from flask_cors import CORS

import agent as tri_agent
import kernel_bridge as kb

app = Flask(__name__, static_folder="static")
app.secret_key = os.urandom(24)
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024  
CORS(app, supports_credentials=True)

_histories: dict = {}

UPLOAD_DIR = Path(tempfile.gettempdir()) / "tri_uploads"
UPLOAD_DIR.mkdir(exist_ok=True)

USERS_DIR  = Path("./users")
USERS_DIR.mkdir(exist_ok=True)
USERS_FILE = USERS_DIR / "users.json"

CHAT_HISTORY_DIR = Path("./chat_history")
CHAT_HISTORY_DIR.mkdir(exist_ok=True)

ALLOWED_IMAGE_EXT = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp"}

ALLOWED_TEXT_EXT = {
    ".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".html", ".htm", ".css", ".js", ".ts", ".jsx", ".tsx", ".vue", ".svelte",
    ".c", ".h", ".cpp", ".cc", ".cxx", ".hpp", ".hxx",
    ".cs", ".java", ".rs", ".go", ".zig",
    ".py", ".rb", ".php", ".pl", ".lua", ".r", ".m", ".sh", ".bash", ".zsh", ".fish", ".ps1",
    ".sql", ".graphql", ".proto",
    ".swift", ".kt", ".kts", ".dart", ".scala", ".ex", ".exs", ".erl", ".hs", ".ml",
    ".f", ".f90", ".f95", ".asm", ".s", ".makefile", ".cmake",
}

ALLOWED_BINARY_EXT = {
    ".pdf",
    ".docx", ".doc",
    ".xlsx", ".xls",
    ".pptx",
}

def _hash_password(password: str) -> str:
    return hashlib.sha256(password.encode("utf-8")).hexdigest()


def load_users() -> dict:
    if not USERS_FILE.exists():
        return {}
    try:
        with open(USERS_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}


def save_users(users: dict) -> None:
    with open(USERS_FILE, "w", encoding="utf-8") as f:
        json.dump(users, f, ensure_ascii=False, indent=2)


def get_sessions_file(username: str) -> Path:
    safe = "".join(c if c.isalnum() or c in "-_" else "_" for c in username.lower())
    return CHAT_HISTORY_DIR / f"chat_sessions_{safe}.json"

@app.route("/auth/register", methods=["POST"])
def auth_register():
    data     = request.get_json(force=True)
    name     = (data.get("name") or "").strip()
    age      = data.get("age")
    gender   = (data.get("gender") or "").strip()
    username = (data.get("username") or "").strip().lower()
    password = (data.get("password") or "").strip()

    if not all([name, age, gender, username, password]):
        return jsonify({"error": "All fields are required."}), 400
    if len(username) < 3:
        return jsonify({"error": "Username must be at least 3 characters."}), 400
    if len(password) < 4:
        return jsonify({"error": "Password must be at least 4 characters."}), 400
    if not str(age).isdigit() or not (1 <= int(age) <= 120):
        return jsonify({"error": "Please enter a valid age."}), 400

    users = load_users()
    if username in users:
        return jsonify({"error": "Username already taken. Please choose another."}), 409

    users[username] = {
        "name":       name,
        "age":        int(age),
        "gender":     gender,
        "username":   username,
        "password":   _hash_password(password),
        "created_at": time.time(),
    }
    save_users(users)

    session["username"] = username
    session["sid"]      = str(uuid.uuid4())

    pid = kb.get_or_create_pid(username)
    print(f"\n[OS] Registered user '{username}' → PID={pid}")

    return jsonify({"status": "registered", "username": username, "name": name, "pid": pid})


@app.route("/auth/login", methods=["POST"])
def auth_login():
    data     = request.get_json(force=True)
    username = (data.get("username") or "").strip().lower()
    password = (data.get("password") or "").strip()

    if not username or not password:
        return jsonify({"error": "Username and password are required."}), 400

    users = load_users()
    user  = users.get(username)
    if not user or user["password"] != _hash_password(password):
        return jsonify({"error": "Invalid username or password."}), 401

    session["username"] = username
    session["sid"]      = str(uuid.uuid4())

    pid = kb.get_or_create_pid(username)
    print(f"\n[OS] User '{username}' logged in → PID={pid}")
    kb.print_pid_registry()

    return jsonify({
        "status":   "logged_in",
        "username": username,
        "name":     user["name"],
        "age":      user["age"],
        "gender":   user["gender"],
        "pid":      pid,
    })


@app.route("/auth/logout", methods=["POST"])
def auth_logout():
    username = session.pop("username", None)
    session.pop("sid", None)

    if username:
        pid = kb._pid_registry.get(username, -1)
        if pid >= 0:
            print(f"\n[OS] User '{username}' logged out → proc_exit PID={pid}")
            kb.proc_exit(pid)
            kb._pid_registry.pop(username, None)

    return jsonify({"status": "logged_out", "username": username})


@app.route("/auth/me", methods=["GET"])
def auth_me():
    username = session.get("username")
    if not username:
        return jsonify({"authenticated": False}), 401
    users = load_users()
    user  = users.get(username)
    if not user:
        return jsonify({"authenticated": False}), 401
    pid = kb._pid_registry.get(username, -1)
    return jsonify({
        "authenticated": True,
        "username": username,
        "name":     user["name"],
        "age":      user["age"],
        "gender":   user["gender"],
        "pid":      pid,
    })


@app.route("/auth/users", methods=["GET"])
def auth_list_users():
    users = load_users()
    return jsonify({
        "users": [
            {"username": u, "name": v["name"]}
            for u, v in users.items()
        ]
    })


def require_auth():
    username = session.get("username")
    if not username:
        return None, (jsonify({"error": "Not authenticated"}), 401)
    return username, None

def get_history(username: str) -> list:
    if username not in _histories:
        _histories[username] = []
    return _histories[username]


def _allowed_file(filename: str) -> tuple:
    ext  = Path(filename).suffix.lower()
    base = Path(filename).name.lower()
    if base in ("makefile", "dockerfile", "gemfile", "procfile", "rakefile"):
        return True, "text"
    if ext in ALLOWED_IMAGE_EXT:
        return True, "image"
    if ext in ALLOWED_TEXT_EXT:
        return True, "text"
    if ext in ALLOWED_BINARY_EXT:
        return True, "binary"
    return False, "unknown"


def _extract_pdf(path: Path) -> str:
    try:
        import fitz
        doc   = fitz.open(str(path))
        pages = []
        for i, page in enumerate(doc):
            text = page.get_text("text")
            if text.strip():
                pages.append(f"[Page {i + 1}]\n{text.strip()}")
        doc.close()
        return "\n\n".join(pages) if pages else "[PDF contained no extractable text]"
    except ImportError:
        return "[PDF extraction requires PyMuPDF: pip install PyMuPDF]"
    except Exception as e:
        return f"[PDF extraction error: {e}]"


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document
        doc   = Document(str(path))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_texts:
                    parts.append(" | ".join(row_texts))
        return "\n".join(parts) if parts else "[Document contained no readable text]"
    except ImportError:
        return "[DOCX extraction requires python-docx: pip install python-docx]"
    except Exception as e:
        return f"[DOCX extraction error: {e}]"


def _extract_excel(path: Path) -> str:
    ext = path.suffix.lower()
    try:
        if ext == ".xlsx":
            import openpyxl
            wb    = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            lines = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                lines.append(f"=== Sheet: {sheet_name} ===")
                for row in ws.iter_rows(values_only=True):
                    row_vals = [str(v) if v is not None else "" for v in row]
                    if any(v.strip() for v in row_vals):
                        lines.append("\t".join(row_vals))
            wb.close()
            return "\n".join(lines) if lines else "[Spreadsheet contained no data]"
        elif ext == ".xls":
            import xlrd
            wb    = xlrd.open_workbook(str(path))
            lines = []
            for sheet in wb.sheets():
                lines.append(f"=== Sheet: {sheet.name} ===")
                for rx in range(sheet.nrows):
                    row_vals = [str(sheet.cell_value(rx, cx)) for cx in range(sheet.ncols)]
                    if any(v.strip() for v in row_vals):
                        lines.append("\t".join(row_vals))
            return "\n".join(lines) if lines else "[Spreadsheet contained no data]"
    except ImportError as ie:
        pkg = "openpyxl" if ext == ".xlsx" else "xlrd"
        return f"[Excel extraction requires {pkg}: pip install {pkg}]"
    except Exception as e:
        return f"[Excel extraction error: {e}]"
    return "[Unsupported Excel format]"


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs   = Presentation(str(path))
        lines = []
        for i, slide in enumerate(prs.slides):
            lines.append(f"=== Slide {i + 1} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        return "\n".join(lines) if lines else "[Presentation contained no readable text]"
    except ImportError:
        return "[PPTX extraction requires python-pptx: pip install python-pptx]"
    except Exception as e:
        return f"[PPTX extraction error: {e}]"


def _extract_binary(path: Path, ext: str) -> str:
    if ext == ".pdf":
        return _extract_pdf(path)
    elif ext in (".docx", ".doc"):
        return _extract_docx(path)
    elif ext in (".xlsx", ".xls"):
        return _extract_excel(path)
    elif ext == ".pptx":
        return _extract_pptx(path)
    return "[Unsupported binary format]"


def _handle_upload(file) -> dict | None:
    if not file or not file.filename:
        return None
    ok, ftype = _allowed_file(file.filename)
    if not ok:
        return None
    ext      = Path(file.filename).suffix.lower()
    tmp_name = f"{uuid.uuid4().hex}{ext}"
    tmp_path = UPLOAD_DIR / tmp_name
    file.save(str(tmp_path))

    if ftype == "image":
        mime_map  = {".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                     ".png": "image/png",  ".gif":  "image/gif",
                     ".webp": "image/webp", ".bmp": "image/bmp"}
        mime_type = mime_map.get(ext, "image/jpeg")
        with open(tmp_path, "rb") as f:
            b64 = base64.b64encode(f.read()).decode("utf-8")
        return {"type": "image", "filename": file.filename,
                "filepath": str(tmp_path), "base64": b64, "mime_type": mime_type}

    elif ftype == "binary":
        content    = _extract_binary(tmp_path, ext)
        char_limit = 12000
        truncated  = len(content) > char_limit
        label_map  = {".pdf": "PDF document", ".docx": "Word document", ".doc": "Word document",
                      ".xlsx": "Excel spreadsheet", ".xls": "Excel spreadsheet",
                      ".pptx": "PowerPoint presentation"}
        return {"type": "text", "filename": file.filename, "filepath": str(tmp_path),
                "content": content[:char_limit], "truncated": truncated,
                "doc_type": label_map.get(ext, "document")}

    else:
        try:
            content = tmp_path.read_text(encoding="utf-8", errors="replace")
        except Exception:
            content = tmp_path.read_text(encoding="latin-1", errors="replace")
        char_limit = 8000
        return {"type": "text", "filename": file.filename, "filepath": str(tmp_path),
                "content": content[:char_limit], "truncated": len(content) > char_limit}

def load_chat_sessions(username: str) -> list:
    f = get_sessions_file(username)
    if not f.exists():
        return []
    try:
        with open(f, "r", encoding="utf-8") as fp:
            data = json.load(fp)
        return data if isinstance(data, list) else []
    except Exception as e:
        print(f"[Sessions] Failed to load for {username}: {e}")
        return []


def save_chat_sessions(username: str, sessions: list) -> None:
    try:
        with open(get_sessions_file(username), "w", encoding="utf-8") as fp:
            json.dump(sessions, fp, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f"[Sessions] Failed to save for {username}: {e}")


def _os_agent_call(username: str, msg: str, history: list,
                   file_context: dict | None = None) -> dict:
    """
    Full OS-aware agent call.  Per-turn kernel workflow:

      1. Ensure user has a kernel PID  (proc_fork if new)
      2. proc_run   →  RUNNING  (user turn starts)
      3. proc_wait  →  WAITING  (LLM/tool I/O burst)
      4. mem_alloc  →  page current turn into kernel RAM
      5. kernel_manage_context → LRU eviction if RAM full
      6. Agent runs (AI side untouched)
      7. proc_run   →  RUNNING  (response ready)
      8. mem_alloc  →  page response into kernel RAM
      9. kernel_stats printed as OS summary
    """
    pid = kb.get_or_create_pid(username)

    kb.proc_run(pid)

    kb.proc_wait(pid)

    turn_id = len(history)
    user_content = msg or (f"[file:{file_context.get('filename', 'upload')}]"
                            if file_context else "[empty]")
    kb.mem_alloc(pid, turn_id, user_content[:400])

    managed_history = kb.kernel_manage_context(
        history, pid,
        max_active=10,
        save_fn=lambda text, source="offloaded_turn":
            tri_agent.save_to_memory(text, source=source, username=username),
    )

    result = tri_agent.agent(msg, managed_history,
                              file_context=file_context, username=username)

    kb.proc_run(pid)

    answer_id = turn_id + 1
    kb.mem_alloc(pid, answer_id, result["answer"][:400])

    kb.kernel_stats()

    return result

@app.route("/")
def index():
    return send_from_directory(".", "index.html")


@app.route("/chat", methods=["POST"])
def chat():
    username, err = require_auth()
    if err:
        return err

    if request.content_type and "multipart" in request.content_type:
        msg  = (request.form.get("message") or "").strip()
        file = request.files.get("file")
    else:
        data = request.get_json(force=True)
        msg  = (data.get("message") or "").strip()
        file = None

    if not msg and not file:
        return jsonify({"error": "Empty message"}), 400

    history      = get_history(username)
    start        = time.time()
    file_context = _handle_upload(file) if file else None

    try:
        result = _os_agent_call(username, msg, history, file_context=file_context)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

    pid = kb._pid_registry.get(username, -1)
    return jsonify({
        "answer":  result["answer"],
        "steps":   result["steps"],
        "tools":   [t["tool"] for t in result["tool_calls"]],
        "elapsed": round(time.time() - start, 1),
        "status":  result["status"],
        "os_pid":  pid,
    })


@app.route("/paste-image", methods=["POST"])
def paste_image():
    username, err = require_auth()
    if err:
        return err

    data     = request.get_json(force=True)
    data_url = data.get("dataUrl", "")
    if not data_url.startswith("data:image/"):
        return jsonify({"error": "Invalid image data"}), 400

    try:
        header, b64_data = data_url.split(",", 1)
        mime_type = header.split(":")[1].split(";")[0]
        ext_map   = {"image/png": ".png", "image/jpeg": ".jpg", "image/gif": ".gif",
                     "image/webp": ".webp", "image/bmp": ".bmp"}
        ext       = ext_map.get(mime_type, ".png")
        raw       = base64.b64decode(b64_data)
        tmp_name  = f"{uuid.uuid4().hex}{ext}"
        tmp_path  = UPLOAD_DIR / tmp_name
        tmp_path.write_bytes(raw)
    except Exception as e:
        return jsonify({"error": f"Image decode error: {e}"}), 400

    file_context = {"type": "image", "filename": f"pasted-image{ext}",
                    "filepath": str(tmp_path), "base64": b64_data, "mime_type": mime_type}
    token = uuid.uuid4().hex
    _paste_store[token] = file_context
    return jsonify({"token": token})


_paste_store: dict = {}


@app.route("/stream", methods=["GET", "POST"])
def stream():
    username, err = require_auth()
    if err:
        return err

    file        = None
    paste_token = None

    if request.method == "POST":
        ct = request.content_type or ""
        if "multipart" in ct:
            msg  = (request.form.get("message") or "").strip()
            file = request.files.get("file")
        else:
            body        = request.get_json(force=True) or {}
            msg         = (body.get("message") or "").strip()
            paste_token = body.get("pasteToken")
    else:
        msg = (request.args.get("message") or "").strip()

    if not msg and not file and not paste_token:
        return jsonify({"error": "Empty message"}), 400

    history = get_history(username)

    if paste_token and paste_token in _paste_store:
        file_context = _paste_store.pop(paste_token)
    elif file:
        file_context = _handle_upload(file)
    else:
        file_context = None

    pid     = kb.get_or_create_pid(username)
    turn_id = len(history)
    user_content = msg or (f"[file:{file_context.get('filename','upload')}]"
                            if file_context else "[empty]")
    kb.proc_run(pid)
    kb.proc_wait(pid)
    kb.mem_alloc(pid, turn_id, user_content[:400])
    managed_history = kb.kernel_manage_context(
        history, pid, max_active=10,
        save_fn=lambda text, source="offloaded_turn":
            tri_agent.save_to_memory(text, source=source, username=username),
    )

    def generate():
        final_answer = ""
        try:
            for event in tri_agent.agent_streaming(msg, managed_history,
                                                    file_context=file_context,
                                                    username=username):
                payload = json.dumps(event)
                yield f"data: {payload}\n\n"
                if event["event"] == "final":
                    final_answer = event["data"].get("answer", "")
        except Exception as e:
            err_payload = json.dumps({"event": "error", "data": {"message": str(e)}})
            yield f"data: {err_payload}\n\n"

        if final_answer:
            kb.proc_run(pid)
            kb.mem_alloc(pid, turn_id + 1, final_answer[:400])
            kb.kernel_stats()

        yield 'data: {"event": "done"}\n\n'

    return Response(
        generate(),
        mimetype="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )

@app.route("/memory", methods=["GET"])
def memory_info():
    username = session.get("username")
    count    = tri_agent.get_memory_count(username)
    samples  = []
    for doc, meta in tri_agent.get_memory_samples(limit=5, username=username):
        samples.append({
            "text":   doc[:120] + ("..." if len(doc) > 120 else ""),
            "source": meta.get("source", "?"),
        })

    pid = kb._pid_registry.get(username, -1)
    if pid >= 0:
        kb.proc_status(pid)

    return jsonify({"chunks": count, "samples": samples})


@app.route("/clear", methods=["POST"])
def clear_history():
    username = session.get("username")
    if username and username in _histories:
        _histories[username] = []
    return jsonify({"status": "cleared"})


@app.route("/sessions", methods=["GET"])
def get_sessions():
    username, err = require_auth()
    if err:
        return err
    sessions = load_chat_sessions(username)
    return jsonify({"sessions": sessions})


@app.route("/sessions", methods=["POST"])
def save_sessions():
    username, err = require_auth()
    if err:
        return err
    data     = request.get_json(force=True)
    sessions = data.get("sessions", [])
    if not isinstance(sessions, list):
        return jsonify({"error": "sessions must be a list"}), 400
    save_chat_sessions(username, sessions)
    return jsonify({"status": "saved", "count": len(sessions)})


@app.route("/sessions/<session_id>", methods=["DELETE"])
def delete_session(session_id: str):
    username, err = require_auth()
    if err:
        return err
    sessions     = load_chat_sessions(username)
    original_len = len(sessions)
    sessions     = [s for s in sessions if s.get("id") != session_id]
    save_chat_sessions(username, sessions)
    deleted = original_len - len(sessions)
    return jsonify({"status": "deleted", "deleted": deleted})


@app.route("/kernel/status", methods=["GET"])
def kernel_status():
    """Return live kernel health + RAM/disk stats as JSON."""
    online = kb.is_kernel_online()
    stat   = kb._mem_status_raw()
    pids   = {u: p for u, p in kb._pid_registry.items()}

    print("\n[OS] /kernel/status requested")
    kb.kernel_stats()

    return jsonify({
        "kernel_online":  online,
        "mode":           "live" if online else "simulation",
        "ram_used":       stat["ram_used"],
        "ram_capacity":   stat["ram_cap"],
        "disk_pages":     stat["disk_pages"],
        "active_pids":    pids,
        "syscall_count":  kb._call_counter,
    })


@app.route("/kernel/process/<int:pid>", methods=["GET"])
def kernel_process(pid: int):
    """Return live PCB data for a given PID."""
    print(f"\n[OS] /kernel/process/{pid} requested")
    res  = kb.proc_status(pid)
    data = kb._parse_data(res)
    return jsonify(data if data else {"error": "pid_not_found"})


@app.route("/kernel/mem/search", methods=["GET"])
def kernel_mem_search():
    """Perform a cosine similarity search on VectorDB disk."""
    query = request.args.get("q", "")
    top_k = int(request.args.get("top_k", 3))
    if not query:
        return jsonify({"error": "query param 'q' required"}), 400
    res = kb.mem_search(query, top_k)
    return jsonify(res)


@app.route("/kernel/demo", methods=["POST"])
def kernel_demo():
    """
    Run the full syscall demo (same as python kernel_bridge.py)
    in the background and return immediately.
    This is great for showing the evaluator all OS features at once.
    """
    import threading

    def _run_demo():
        print("\n" + "=" * 62)
        print("  [KERNEL DEMO]  Running all syscall tests...")
        print("=" * 62)

        pid = kb.proc_fork("demo_user", priority=2)
        kb.proc_run(pid)

        for i in range(12):
            kb.mem_alloc(pid, i, f"Turn {i}: user asked about topic_{i} concepts.")

        kb.mem_read(pid, 11, query="topic_11")

        kb.mem_read(pid, 0, query="topic_0 explanation")

        kb.mem_search("topic_3 concepts", top_k=3)

        kb.fs_write(pid, "/tmp/llmos_demo.txt", "LLM-OS demo write\n")
        kb.fs_read(pid, "/tmp/llmos_demo.txt")

        kb.proc_status(pid)

        kb.kernel_stats()

        kb.proc_exit(pid)
        print("\n[KERNEL DEMO] Complete.\n")

    t = threading.Thread(target=_run_demo, daemon=True)
    t.start()
    return jsonify({"status": "demo_started",
                    "message": "Watch your terminal for OS output."})


def _startup_banner():
    print("\n" + "=" * 62)
    print("  TRI Chatbot  v2  — OS Integrated")
    print("  Open:  http://localhost:5000")
    print("  Stop:  Ctrl + C")
    print(f"  Users file:       {USERS_FILE.resolve()}")
    print(f"  Chat history dir: {CHAT_HISTORY_DIR.resolve()}")
    print()
    print("  OS KERNEL ENDPOINTS:")
    print("    GET  /kernel/status              Live RAM/disk stats")
    print("    GET  /kernel/process/<pid>        PCB info")
    print("    GET  /kernel/mem/search?q=<query> VectorDB cosine search")
    print("    POST /kernel/demo                 Run full syscall demo")
    print()
    print("  TO USE REAL C++ KERNEL:")
    print("    cd kernel && g++ -std=c++17 -pthread -o kernel_server kernel_server.cpp")
    print("    ./kernel_server   (keep running in a separate terminal)")
    print("=" * 62 + "\n")

    online = kb.is_kernel_online()
    print(f"  Kernel status: {'CONNECTED (C++ live)' if online else 'OFFLINE (Python simulation active)'}")
    print()


if __name__ == "__main__":
    _startup_banner()
    app.run(host="0.0.0.0", port=5000, debug=False, threaded=True)